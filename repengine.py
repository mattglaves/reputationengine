#!/usr/bin/python3

import sys, os, json, re
#import docx2txt
from azurecreds import *
from docx2txt import *
from io import StringIO
from boxsdk import JWTAuth
from boxsdk import Client
from pptx import Presentation
from buzzwords import buzzwords
from pdfminer.pdfinterp import PDFResourceManager, PDFPageInterpreter
from pdfminer.converter import TextConverter
from pdfminer.layout import LAParams
from pdfminer.pdfpage import PDFPage
from azure.cognitiveservices.language.spellcheck import SpellCheckAPI
from azure.cognitiveservices.language.textanalytics import TextAnalyticsAPI
from msrest.authentication import CognitiveServicesCredentials

attachment_directory = "files/"
file_types = ('pdf','txt','docx','pptx')

os.chdir('/var/www/repengine')

sdk = JWTAuth.from_settings_file('./38487735_v17qtrdu_config.json')
Box = Client(sdk)

def read_text (file_name):

    notes = ''

    if file_name.endswith('.pdf'):
        rsrcmgr = PDFResourceManager()
        retstr = StringIO()
        laparams = LAParams()
        device = TextConverter(rsrcmgr, retstr, codec='utf-8', laparams=laparams)
        fp = open(file_name, 'rb')
        interpreter = PDFPageInterpreter(rsrcmgr, device)
        pagenos=set()

        for page in PDFPage.get_pages(fp, pagenos, maxpages=0, password="",caching=True, check_extractable=True):
                interpreter.process_page(page)

        document_text = retstr.getvalue()
        fp.close()
        device.close()
        retstr.close()
    elif file_name.endswith('.docx'):
        (document_text,notes) = process(file_name)
    elif file_name.endswith('.txt'):
        with open(file_name, 'r') as txt_file:
                document_text = txt_file.read()
    elif file_name.endswith('.pptx'):
        prs = Presentation(file_name)
        slide_txt = ''
        for slide in prs.slides:
            if slide.has_notes_slide:
                notes += re.sub(r'https?:\/\/.*[\r\n]*', '', slide.notes_slide.notes_text_frame.text, flags=re.MULTILINE)
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                            slide_txt += run.text + ' '
            document_text = slide_txt

    document_text = document_text.replace('\n', ' ').replace('\r', '')
    return (document_text,notes)

def check_language (document_text):

    documents = [{'id': 0,'text': document_text[:4999]}]
    analytics_client = TextAnalyticsAPI("westcentralus", CognitiveServicesCredentials(AZURE_KEY_ANALYTICS))
    language_response = analytics_client.detect_language(documents=documents)
    return (language_response.documents[0].detected_languages[0].name)

def check_spelling (document_text):

    errors = 0
    spellcheck_client = SpellCheckAPI(CognitiveServicesCredentials(AZURE_KEY_SPELLING))
    result = spellcheck_client.spell_checker(document_text[:9000], mode="proof")

    for x in range(0, len(result.flagged_tokens)):
        if result.flagged_tokens[x].token.isalpha():
            if result.flagged_tokens[x].suggestions[0].score > .85:
                errors += 1
    return (errors)

def check_sentiment (document_text):
    
    documents = [{'language': 'en', 'id': 0,'text': document_text[:500]}]
    sentiment_client = TextAnalyticsAPI("westcentralus", CognitiveServicesCredentials(AZURE_KEY_ANALYTICS))
    sentiment_response = sentiment_client.sentiment(documents=documents)
    return (sentiment_response.documents[0].score)

def maximum_hype (document_text):

    hype_score = 0
    for i in range(len(buzzwords)):
        if document_text.find(buzzwords[i]) > 0:
            #print (buzzwords[i])
            hype_score += 1
    return (hype_score)

def sharing_restrictions (file_id,BoxUser,spelling_errors, sentiment_score, hype_score, notes):

    found = 0
    skills_card = {
        "cards" : [
        {
            "skill_card_title" : {
                "code" : "skills_image_text",
                "message" : "External Sharing Prohibited "
            },
            "invocation" : {
                "type" : "skill_invocation",
                "id" : "6e7ad73c-1a56-48c6-82c0-3563abe25c38_297548314"
            },        
            "skill" : {
            "type" : "service",
            "id" : "329"
            },
            "status" : {},
            "skill_card_type" : "transcript",
            "type" : "skill_card",
            "entries" : [
            {
                "type" : "text",
                "text" : "This file can not be shared until the following items are resolved:"
            }
        ]
        }]}
    
    if spelling_errors > 1:
        skills_card["cards"][0]["entries"].append({ "type" : "text", "text" : "- Spelling and/or Grammar Errors" })
    if sentiment_score < .4:
        skills_card["cards"][0]["entries"].append({ "type" : "text", "text" : "- Negative Document Context - Improve tone" })
    if hype_score > 3:
        skills_card["cards"][0]["entries"].append({ "type" : "text", "text" : "- Buzzword Overload - Reduce synergy" })
    if notes > 1:
        skills_card["cards"][0]["entries"].append({ "type" : "text", "text" : "- Document Comments and/or Speaker & Editor Notes not Removed" })

    if len(skills_card["cards"][0]["entries"]) > 1:
        metadata = Box.as_user(BoxUser).file(file_id).metadata('global', 'boxSkillsCards')
        metadata.create(skills_card)
        metadata = Box.as_user(BoxUser).file(file_id).metadata('enterprise', 'securityClassification-6VMVochwUWo')
        metadata.create({"Box__Security__Classification__Key" : "Internal"})

def remove_restrictions (file_id,file_owner):
    print ("Sup - I'm totally not implemented! ¯\_(ツ)_/¯")

def main():


    if (len(sys.argv) == 4):
        file_id = sys.argv[1]
        file_name = attachment_directory+sys.argv[2].lower()
        file_owner = sys.argv[3]

        BoxUser = Box.user(user_id=file_owner)

        if file_name.endswith(file_types):
            with open(file_name, 'wb') as open_file:
                Box.as_user(BoxUser).file(file_id).download_to(open_file)

            (document_text, notes) = read_text(file_name) 
            if (check_language(document_text) == 'English'):
                spelling_errors = check_spelling(document_text)
                sentiment_score = check_sentiment(document_text)
                hype_score = maximum_hype(document_text)

                if spelling_errors > 1 or sentiment_score < .3 or hype_score > 3 or len(notes) > 2:
                    sharing_restrictions (file_id,BoxUser,spelling_errors, sentiment_score, hype_score, len(notes))
                #else:
                    #remove_restrictions (file_id,BoxUser)

main()
